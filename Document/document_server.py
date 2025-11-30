#!/usr/bin/env python3

"""
Simple Document Generator MCP Server - Generate Word, Excel, and PowerPoint documents with templates
"""

import os
import sys
import logging
import json
from datetime import datetime, timezone
import httpx
from mcp.server.fastmcp import FastMCP

# Document libraries
from docx import Document
import docx
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.enum.table import WD_TABLE_ALIGNMENT
import openpyxl
from pptx import Presentation
from reportlab.pdfgen import canvas
from reportlab.lib.pagesizes import letter
from docx.enum.text import WD_ALIGN_PARAGRAPH

# Configure logging to stderr
logging.basicConfig(
    level=logging.INFO,
    format='%(asctime)s - %(name)s - %(levelname)s - %(message)s',
    stream=sys.stderr
)
logger = logging.getLogger("document-server")

# Initialize MCP server - NO PROMPT PARAMETER!
mcp = FastMCP("document")

# Configuration
# Add any API keys, URLs, or configuration here
# API_TOKEN = os.environ.get("DOCUMENT_API_TOKEN", "")

# === UTILITY FUNCTIONS ===

def ensure_output_dir():
    """Ensure output directory exists."""
    os.makedirs("/app/outputs", exist_ok=True)

def parse_markdown_table(doc, table_md, insert_after=None):
    """Parse markdown table and add to docx document."""
    lines = table_md.strip().split('\n')
    if len(lines) < 2:
        return
    
    # Assume first line is header, second is separator
    headers = [h.strip() for h in lines[0].split('|')[1:-1]]
    data_rows = []
    
    for line in lines[2:]:
        if line.strip() and not line.startswith('|---'):
            cells = [c.strip() for c in line.split('|')[1:-1]]
            data_rows.append(cells)
    
    # Add table at end first
    table = doc.add_table(rows=1, cols=len(headers))
    table.style = 'Table Grid'
    table.alignment = WD_TABLE_ALIGNMENT.CENTER  # Center the table
    
    hdr_cells = table.rows[0].cells
    for i, h in enumerate(headers):
        hdr_cells[i].text = h
        # Make header bold and center
        for paragraph in hdr_cells[i].paragraphs:
            for run in paragraph.runs:
                run.bold = True
            paragraph.alignment = WD_ALIGN_PARAGRAPH.CENTER
    
    for row_data in data_rows:
        row_cells = table.add_row().cells
        for i, cell in enumerate(row_data):
            if i < len(row_cells):
                row_cells[i].text = cell
                # Left align data cells
                for paragraph in row_cells[i].paragraphs:
                    paragraph.alignment = WD_ALIGN_PARAGRAPH.LEFT
    
    # If insert_after, move table after that paragraph
    if insert_after:
        insert_after._element.addnext(table._element)
        # Remove from current position
        doc._body._element.remove(table._element)

def replace_placeholders_in_doc(doc, placeholders):
    """Replace placeholders in a Word document, including images and tables."""
    for paragraph in doc.paragraphs:
        for key, value in placeholders.items():
            # Check for both {key} and [{key}]
            plain_key = key
            bracketed_key = f'[{key}]'
            if plain_key in paragraph.text:
                target = plain_key
            elif bracketed_key in paragraph.text:
                target = bracketed_key
            else:
                continue
            
            if key.startswith('{image'):
                # Handle image placeholder
                image_path = value
                if os.path.exists(image_path):
                    paragraph.clear()
                    run = paragraph.add_run()
                    run.add_picture(image_path, width=docx.shared.Inches(6))
                    paragraph.alignment = WD_ALIGN_PARAGRAPH.CENTER  # Center the image
                else:
                    paragraph.text = paragraph.text.replace(target, f"[Image not found: {image_path}]")
            elif key == '{findings}' and '|' in str(value):
                # Handle markdown table
                paragraph.clear()
                parse_markdown_table(doc, str(value))
            else:
                # Handle text placeholders with multi-paragraph support
                parts = str(value).split('\n\n')
                if len(parts) == 1:
                    paragraph.text = parts[0]
                    if '**' in parts[0] or '*' in parts[0]:
                        parse_markdown_to_runs(paragraph, parts[0])
                    paragraph.alignment = WD_ALIGN_PARAGRAPH.LEFT
                else:
                    # Multi-paragraph text
                    paragraph.text = parts[0]
                    if '**' in parts[0] or '*' in parts[0]:
                        parse_markdown_to_runs(paragraph, parts[0])
                    paragraph.alignment = WD_ALIGN_PARAGRAPH.LEFT
                    # Insert subsequent parts as new paragraphs
                    for part in parts[1:]:
                        new_p_element = docx.oxml.shared.OxmlElement("w:p")
                        paragraph._element.addnext(new_p_element)
                        new_p = docx.text.paragraph.Paragraph(new_p_element, doc)
                        new_p.text = part
                        if '**' in part or '*' in part:
                            parse_markdown_to_runs(new_p, part)
                        new_p.alignment = WD_ALIGN_PARAGRAPH.LEFT
                        paragraph = new_p
    
    # Also replace in tables
    for table in doc.tables:
        for row in table.rows:
            for cell in row.cells:
                for paragraph in cell.paragraphs:
                    for key, value in placeholders.items():
                        plain_key = key
                        bracketed_key = f'[{key}]'
                        if plain_key in paragraph.text:
                            target = plain_key
                        elif bracketed_key in paragraph.text:
                            target = bracketed_key
                        else:
                            continue
                        
                        if key.startswith('{image'):
                            # Handle image placeholder
                            image_path = value
                            if os.path.exists(image_path):
                                paragraph.clear()
                                run = paragraph.add_run()
                                run.add_picture(image_path, width=docx.shared.Inches(6))
                                paragraph.alignment = WD_ALIGN_PARAGRAPH.CENTER
                            else:
                                paragraph.text = paragraph.text.replace(target, f"[Image not found: {image_path}]")
                        elif key == '{findings}' and '|' in str(value):
                            # Handle markdown table in table cell - unlikely but handle
                            paragraph.clear()
                            parse_markdown_table(doc, str(value), insert_after=paragraph)
                        else:
                            # Handle text placeholders in table cells
                            parts = str(value).split('\n\n')
                            if len(parts) == 1:
                                paragraph.text = parts[0]
                                if '**' in parts[0] or '*' in parts[0]:
                                    parse_markdown_to_runs(paragraph, parts[0])
                                paragraph.alignment = WD_ALIGN_PARAGRAPH.LEFT
                            else:
                                # Multi-paragraph in table cell - replace with first part, ignore rest for simplicity
                                paragraph.text = parts[0]
                                if '**' in parts[0] or '*' in parts[0]:
                                    parse_markdown_to_runs(paragraph, parts[0])
                                paragraph.alignment = WD_ALIGN_PARAGRAPH.LEFT

def parse_markdown_to_runs(para, content):
    """Parse content with simple MD-like formatting (**bold**, *italic*) and add to paragraph."""
    import re
    lines = content.split('\n')
    for line in lines:
        if not line.strip():
            para.add_run('\n')
            continue
        # Split by ** and *
        parts = re.split(r'(\*\*.*?\*\*|\*.*?\*)', line)
        for part in parts:
            run = para.add_run()
            if part.startswith('**') and part.endswith('**'):
                run.text = part[2:-2]
                run.bold = True
            elif part.startswith('*') and part.endswith('*'):
                run.text = part[1:-1]
                run.italic = True
            else:
                run.text = part
        para.add_run('\n')  # Add newline after each line

# === MCP TOOLS ===

@mcp.tool()
async def document_generate_word(content: str = "", template: str = "", output_name: str = "", placeholders: str = "") -> str:
    """Generate a Word document with optional template, content, and placeholders to fill."""
    logger.info(f"Generating Word document with content length {len(content)}, template: {template}, output: {output_name}, placeholders: {placeholders}")
    
    if not content.strip() and not placeholders.strip():
        return "❌ Error: Content or placeholders are required"
    
    if not output_name.strip():
        output_name = "document.docx"
    
    try:
        ensure_output_dir()
        output_path = f"/app/outputs/{output_name}"
        
        if template.strip():
            template_path = f"/app/templates/{template}"
            if not os.path.exists(template_path):
                return f"❌ Error: Template not found: {template_path}"
            # Handle .dotx templates by copying to temp .docx
            if template.lower().endswith('.dotx'):
                import shutil
                temp_docx = f"/tmp/{os.path.basename(template_path).replace('.dotx', '.docx')}"
                shutil.copy2(template_path, temp_docx)
                doc = Document(temp_docx)
            else:
                doc = Document(template_path)
        else:
            doc = Document()
        
        # Parse placeholders if provided
        placeholder_dict = {}
        if placeholders.strip():
            try:
                placeholder_dict = json.loads(placeholders)
            except json.JSONDecodeError:
                return "❌ Error: Invalid JSON for placeholders"
        
        # Replace placeholders in the document
        if placeholder_dict:
            replace_placeholders_in_doc(doc, placeholder_dict)
        
        # Add content with formatting
        if content.strip():
            parse_content_to_paragraph(doc, content)
        
        doc.save(output_path)
        
        return f"✅ Word document generated successfully: {output_path}"
    except Exception as e:
        logger.error(f"Error generating Word document: {e}")
        return f"❌ Error: {str(e)}"

@mcp.tool()
async def document_generate_excel(data: str = "", template: str = "", output_name: str = "") -> str:
    """Generate an Excel document with data and optional template."""
    logger.info(f"Generating Excel document with data length {len(data)}, template: {template}, output: {output_name}")
    
    if not data.strip():
        return "❌ Error: Data is required (JSON format)"
    
    if not output_name.strip():
        output_name = "document.xlsx"
    
    try:
        ensure_output_dir()
        output_path = f"/app/outputs/{output_name}"
        
        if template.strip():
            template_path = f"/app/templates/{template}"
            if not os.path.exists(template_path):
                return f"❌ Error: Template not found: {template_path}"
            wb = openpyxl.load_workbook(template_path)
        else:
            wb = openpyxl.Workbook()
        
        ws = wb.active
        parsed_data = json.loads(data)
        
        if isinstance(parsed_data, list):
            for row in parsed_data:
                if isinstance(row, list):
                    ws.append(row)
                else:
                    ws.append([row])
        else:
            return "❌ Error: Data must be a JSON array of arrays"
        
        wb.save(output_path)
        
        return f"✅ Excel document generated successfully: {output_path}"
    except json.JSONDecodeError:
        return "❌ Error: Invalid JSON data"
    except Exception as e:
        logger.error(f"Error generating Excel document: {e}")
        return f"❌ Error: {str(e)}"

@mcp.tool()
async def document_generate_powerpoint(slides: str = "", template: str = "", output_name: str = "") -> str:
    """Generate a PowerPoint document with slides and optional template."""
    logger.info(f"Generating PowerPoint document with slides length {len(slides)}, template: {template}, output: {output_name}")
    
    if not slides.strip():
        return "❌ Error: Slides data is required (JSON format)"
    
    if not output_name.strip():
        output_name = "document.pptx"
    
    try:
        ensure_output_dir()
        output_path = f"/app/outputs/{output_name}"
        
        if template.strip():
            template_path = f"/app/templates/{template}"
            if not os.path.exists(template_path):
                return f"❌ Error: Template not found: {template_path}"
            prs = Presentation(template_path)
        else:
            prs = Presentation()
        
        parsed_slides = json.loads(slides)
        
        if isinstance(parsed_slides, list):
            for slide_data in parsed_slides:
                if isinstance(slide_data, dict):
                    slide = prs.slides.add_slide(prs.slide_layouts[1])
                    slide.shapes.title.text = slide_data.get("title", "")
                    if len(slide.placeholders) > 1:
                        slide.placeholders[1].text = slide_data.get("content", "")
                else:
                    return "❌ Error: Each slide must be a JSON object with 'title' and 'content'"
        else:
            return "❌ Error: Slides must be a JSON array of objects"
        
        prs.save(output_path)
        
        return f"✅ PowerPoint document generated successfully: {output_path}"
    except json.JSONDecodeError:
        return "❌ Error: Invalid JSON slides data"
    except Exception as e:
        logger.error(f"Error generating PowerPoint document: {e}")
        return f"❌ Error: {str(e)}"

@mcp.tool()
async def document_generate_pdf(content: str = "", template: str = "", output_name: str = "") -> str:
    """Generate a PDF document with text content."""
    logger.info(f"Generating PDF document with content length {len(content)}, template: {template}, output: {output_name}")
    
    if not content.strip():
        return "❌ Error: Content is required"
    
    if not output_name.strip():
        output_name = "document.pdf"
    
    try:
        ensure_output_dir()
        output_path = f"/app/outputs/{output_name}"
        
        # For simplicity, ignore template for now, as PDF templates are more complex
        c = canvas.Canvas(output_path, pagesize=letter)
        width, height = letter
        
        # Split content into lines and draw
        lines = content.split('\n')
        y = height - 50  # Start near top
        for line in lines:
            if y < 50:  # New page if near bottom
                c.showPage()
                y = height - 50
            c.drawString(50, y, line)
            y -= 15  # Line spacing
        
        c.save()
        
        return f"✅ PDF document generated successfully: {output_path}"
    except Exception as e:
        logger.error(f"Error generating PDF document: {e}")
        return f"❌ Error: {str(e)}"

# === SERVER STARTUP ===

if __name__ == "__main__":
    logger.info("Starting Document Generator MCP server...")
    
    # Add any startup checks
    # if not API_TOKEN:
    # logger.warning("DOCUMENT_API_TOKEN not set")
    
    try:
        mcp.run(transport='stdio')
    except Exception as e:
        logger.error(f"Server error: {e}", exc_info=True)
        sys.exit(1)